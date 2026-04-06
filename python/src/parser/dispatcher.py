"""多格式文档解析分发器 — 根据文件后缀自动选择提取器"""

from __future__ import annotations

import csv
import json
import re
import xml.etree.ElementTree as ET
from io import StringIO
from pathlib import Path
from typing import Callable

from src.parser.extractor import DocumentContent, PageContent, extract_pages


# ---------------------------------------------------------------------------
# 支持的文件格式 (后缀 → 格式名)
# ---------------------------------------------------------------------------

SUPPORTED_EXTENSIONS: dict[str, str] = {
    ".pdf": "PDF",
    ".txt": "纯文本",
    ".md": "Markdown",
    ".log": "日志",
    ".docx": "Word",
    ".doc": "Word 旧版",
    ".html": "HTML",
    ".htm": "HTML",
    ".epub": "EPUB",
    ".rtf": "RTF",
    ".tex": "LaTeX",
    ".csv": "CSV",
    ".pptx": "PowerPoint",
    ".xlsx": "Excel",
    ".srt": "字幕",
    ".json": "JSON",
    ".xml": "XML",
}


# ---------------------------------------------------------------------------
# 提取器注册表
# ---------------------------------------------------------------------------

_Registry = dict[str, Callable[[Path], DocumentContent]]
_REGISTRY: _Registry = {}


def _register(*exts: str):
    """装饰器：将提取函数注册到指定后缀"""
    def decorator(fn: Callable[[Path], DocumentContent]):
        for ext in exts:
            _REGISTRY[ext.lower()] = fn
        return fn
    return decorator


# ---------------------------------------------------------------------------
# 纯文本类 (无额外依赖)
# ---------------------------------------------------------------------------

def _read_text(path: Path) -> str:
    """读取文本文件，UTF-8 优先，自动检测编码"""
    raw = path.read_bytes()
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        try:
            from charset_normalizer import from_bytes
            result = from_bytes(raw)
            best = result.best()
            if best:
                return str(best)
        except ImportError:
            pass
        return raw.decode("latin-1")


@_register(".txt", ".md", ".log")
def _extract_txt(path: Path) -> DocumentContent:
    text = _read_text(path)
    return DocumentContent(
        pages=[PageContent(page_num=1, text=text, width=0, height=0)],
        source_path=str(path),
    )


@_register(".csv")
def _extract_csv(path: Path) -> DocumentContent:
    text = _read_text(path)
    reader = csv.reader(StringIO(text))
    rows = [" | ".join(row) for row in reader]
    content = "\n".join(rows)
    return DocumentContent(
        pages=[PageContent(page_num=1, text=content, width=0, height=0)],
        source_path=str(path),
    )


@_register(".json")
def _extract_json(path: Path) -> DocumentContent:
    text = _read_text(path)
    try:
        data = json.loads(text)
    except json.JSONDecodeError as e:
        raise ValueError(f"JSON 文件格式错误: {e}") from e
    lines: list[str] = []

    def _walk(obj, prefix: str = ""):
        if isinstance(obj, str):
            if obj.strip():
                label = f"{prefix}: " if prefix else ""
                lines.append(f"{label}{obj}")
        elif isinstance(obj, dict):
            for k, v in obj.items():
                child_prefix = f"{prefix}.{k}" if prefix else k
                _walk(v, child_prefix)
        elif isinstance(obj, list):
            for i, item in enumerate(obj):
                _walk(item, f"{prefix}[{i}]" if prefix else f"[{i}]")

    _walk(data)
    content = "\n".join(lines)
    return DocumentContent(
        pages=[PageContent(page_num=1, text=content, width=0, height=0)],
        source_path=str(path),
    )


@_register(".xml")
def _extract_xml(path: Path) -> DocumentContent:
    try:
        tree = ET.parse(path)
    except ET.ParseError as e:
        raise ValueError(f"XML 文件格式错误: {e}") from e
    root = tree.getroot()
    texts: list[str] = []

    def _collect(elem: ET.Element):
        if elem.text and elem.text.strip():
            texts.append(elem.text.strip())
        if elem.tail and elem.tail.strip():
            texts.append(elem.tail.strip())
        for child in elem:
            _collect(child)

    _collect(root)
    content = "\n".join(texts)
    return DocumentContent(
        pages=[PageContent(page_num=1, text=content, width=0, height=0)],
        source_path=str(path),
    )


@_register(".srt")
def _extract_srt(path: Path) -> DocumentContent:
    text = _read_text(path)
    # SRT 格式: 序号 → 时间戳 → 字幕文本(可能多行) → 空行
    blocks = re.split(r"\n\s*\n", text)
    lines: list[str] = []
    for block in blocks:
        block_lines = block.strip().split("\n")
        # 跳过序号行和时间戳行，保留字幕文本
        for line in block_lines[2:]:
            line = line.strip()
            if line:
                lines.append(line)
    content = " ".join(lines)
    return DocumentContent(
        pages=[PageContent(page_num=1, text=content, width=0, height=0)],
        source_path=str(path),
    )


# ---------------------------------------------------------------------------
# PDF (已有 pdfplumber)
# ---------------------------------------------------------------------------

@_register(".pdf")
def _extract_pdf(path: Path) -> DocumentContent:
    return extract_pages(path)


# ---------------------------------------------------------------------------
# Word (.docx) — python-docx
# ---------------------------------------------------------------------------

@_register(".docx")
def _extract_docx(path: Path) -> DocumentContent:
    try:
        from docx import Document
    except ImportError:
        raise ImportError("请安装 python-docx: pip install python-docx")

    doc = Document(str(path))
    parts: list[str] = []

    # 段落
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())

    # 表格
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            parts.append(" | ".join(cells))

    content = "\n\n".join(parts)
    return DocumentContent(
        pages=[PageContent(page_num=1, text=content, width=0, height=0)],
        source_path=str(path),
    )


# ---------------------------------------------------------------------------
# HTML — BeautifulSoup
# ---------------------------------------------------------------------------

@_register(".html", ".htm")
def _extract_html(path: Path) -> DocumentContent:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        raise ImportError("请安装 beautifulsoup4: pip install beautifulsoup4")

    html = _read_text(path)
    soup = BeautifulSoup(html, "html.parser")

    # 移除噪声标签
    for tag in soup.find_all(["script", "style", "nav", "footer", "header", "noscript"]):
        tag.decompose()

    content = soup.get_text(separator="\n", strip=True)
    return DocumentContent(
        pages=[PageContent(page_num=1, text=content, width=0, height=0)],
        source_path=str(path),
    )


# ---------------------------------------------------------------------------
# EPUB — ebooklib
# ---------------------------------------------------------------------------

@_register(".epub")
def _extract_epub(path: Path) -> DocumentContent:
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
    except ImportError:
        raise ImportError("请安装 ebooklib 和 beautifulsoup4")

    book = epub.read_epub(str(path))
    pages: list[PageContent] = []
    idx = 0

    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        html = item.get_content().decode("utf-8", errors="replace")
        soup = BeautifulSoup(html, "html.parser")
        text = soup.get_text(separator="\n", strip=True)
        if text.strip():
            idx += 1
            pages.append(PageContent(page_num=idx, text=text, width=0, height=0))

    if not pages:
        pages.append(PageContent(page_num=1, text="", width=0, height=0))

    return DocumentContent(pages=pages, source_path=str(path))


# ---------------------------------------------------------------------------
# RTF — striprtf
# ---------------------------------------------------------------------------

@_register(".rtf")
def _extract_rtf(path: Path) -> DocumentContent:
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        raise ImportError("请安装 striprtf: pip install striprtf")

    raw = _read_text(path)
    content = rtf_to_text(raw)
    return DocumentContent(
        pages=[PageContent(page_num=1, text=content, width=0, height=0)],
        source_path=str(path),
    )


# ---------------------------------------------------------------------------
# LaTeX — pylatexenc
# ---------------------------------------------------------------------------

@_register(".tex")
def _extract_latex(path: Path) -> DocumentContent:
    try:
        from pylatexenc.latex2text import LatexNodes2Text
    except ImportError:
        raise ImportError("请安装 pylatexenc: pip install pylatexenc")

    raw = _read_text(path)
    content = LatexNodes2Text().latex_to_text(raw)
    return DocumentContent(
        pages=[PageContent(page_num=1, text=content, width=0, height=0)],
        source_path=str(path),
    )


# ---------------------------------------------------------------------------
# PowerPoint (.pptx) — python-pptx
# ---------------------------------------------------------------------------

@_register(".pptx")
def _extract_pptx(path: Path) -> DocumentContent:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("请安装 python-pptx: pip install python-pptx")

    prs = Presentation(str(path))
    pages: list[PageContent] = []

    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text.strip())
        pages.append(PageContent(page_num=i, text="\n".join(texts), width=0, height=0))

    return DocumentContent(pages=pages, source_path=str(path))


# ---------------------------------------------------------------------------
# Excel (.xlsx) — openpyxl
# ---------------------------------------------------------------------------

@_register(".xlsx")
def _extract_xlsx(path: Path) -> DocumentContent:
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ImportError("请安装 openpyxl: pip install openpyxl")

    wb = load_workbook(str(path), read_only=True, data_only=True)
    try:
        pages: list[PageContent] = []

        for i, sheet in enumerate(wb.worksheets, 1):
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None]
                if cells:
                    rows.append(" | ".join(cells))
            pages.append(PageContent(page_num=i, text="\n".join(rows), width=0, height=0))

        return DocumentContent(pages=pages, source_path=str(path))
    finally:
        wb.close()


# ---------------------------------------------------------------------------
# Word 旧版 (.doc) — 提示用户转换
# ---------------------------------------------------------------------------

@_register(".doc")
def _extract_doc(path: Path) -> DocumentContent:
    raise ValueError(
        "旧版 .doc 格式不支持直接解析。"
        "请使用 LibreOffice 转换: soffice --convert-to docx 文件名.doc"
    )


# ---------------------------------------------------------------------------
# 公共 API
# ---------------------------------------------------------------------------

def extract_document(file_path: str | Path) -> DocumentContent:
    """根据文件后缀自动分发到对应提取器

    Args:
        file_path: 文件路径

    Returns:
        DocumentContent 包含提取的文本

    Raises:
        FileNotFoundError: 文件不存在
        ValueError: 不支持的格式
    """
    file_path = Path(file_path)
    if not file_path.exists():
        raise FileNotFoundError(f"文件不存在: {file_path}")

    ext = file_path.suffix.lower()
    extractor = _REGISTRY.get(ext)
    if extractor is None:
        supported = ", ".join(sorted(_REGISTRY.keys()))
        raise ValueError(f"不支持的文件格式: {ext}。支持的格式: {supported}")

    return extractor(file_path)


def get_supported_extensions() -> list[str]:
    """返回所有支持的文件后缀列表"""
    return sorted(_REGISTRY.keys())
